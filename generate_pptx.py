import sys, io, os
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Emu(12192000)
prs.slide_height = Emu(6858000)

DARK_BG = RGBColor(0x1a, 0x1a, 0x2e)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT = RGBColor(0xCF, 0x1F, 0x2E)
TEXT_PRIMARY = RGBColor(0x22, 0x22, 0x22)
TEXT_SECONDARY = RGBColor(0x66, 0x66, 0x66)
TEXT_TERTIARY = RGBColor(0x99, 0x99, 0x99)

LOGO_RED = os.path.join('Gamma 감마_템플릿 세팅 소스', 'BNI로고_transparent.png')
LOGO_WHITE = os.path.join('Gamma 감마_템플릿 세팅 소스', 'BNI로고_white.png')

def add_logo(slide, dark=False):
    path = LOGO_WHITE if dark else LOGO_RED
    if os.path.exists(path):
        slide.shapes.add_picture(path, Emu(10800000), Emu(200000), height=Emu(350000))

def set_bg_dark(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BG

def add_accent_line(slide, top):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(5596000), top, Emu(1000000), Emu(50000))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()

def add_bottom_bar(slide):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(0), Emu(6798000), Emu(12192000), Emu(60000))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()

def add_text(slide, left, top, width, height, text, size=18, bold=False, color=TEXT_PRIMARY, align=PP_ALIGN.CENTER, font_name='Pretendard'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    lines = text.split('\n')
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.color.rgb = color
        p.font.name = font_name
        p.alignment = align
        p.space_after = Pt(4)

def add_site_name(slide, dark=False):
    color = RGBColor(0xCC, 0xCC, 0xDD) if dark else TEXT_PRIMARY
    add_text(slide, Emu(200000), Emu(150000), Emu(3000000), Emu(350000),
             'My PowerTeam', size=11, bold=True, color=color, align=PP_ALIGN.LEFT)

def add_copyright(slide, dark=False):
    color = RGBColor(0x88, 0x88, 0x99) if dark else TEXT_TERTIARY
    add_text(slide, Emu(200000), Emu(6500000), Emu(6000000), Emu(300000),
             'COPYRIGHT© BNI K. JOY CO., LTD. ALL RIGHTS RESERVED.',
             size=8, color=color, align=PP_ALIGN.LEFT)

def add_page_num(slide, num, total, dark=False):
    color = RGBColor(0x88, 0x88, 0x99) if dark else TEXT_TERTIARY
    add_text(slide, Emu(10500000), Emu(6500000), Emu(1500000), Emu(300000),
             f'{num}/{total}', size=9, color=color, align=PP_ALIGN.RIGHT)

slides_data = [
    # === 최종 VERSION (1-11) ===
    {'id': '1', 'bg': 'dark', 'label': 'BNI KOREA POWERTEAM TOOLKIT',
     'title': '무작위 리드에서\n의도적 리퍼럴로',
     'body': '파워팀 툴킷과 함께하는 체계적 리퍼럴 파트너 영입',
     'group': '최종'},

    {'id': '2', 'bg': 'light', 'label': 'THE OLD WAY',
     'title': '지금까지의 방식: 무작위 리드',
     'body': '"Leads are Random" — BNI Global Power Team\n\n• 불확실한 무작위 리드\n• 낮은 신뢰와 일방적 관계\n• 낮은 전환율과 성과 저조',
     'group': '최종'},

    {'id': '3', 'bg': 'light', 'label': 'THE SMART WAY',
     'title': '새로운 방식: 의도적 리퍼럴',
     'body': '"Referrals are Intentional" — BNI Global Power Team\n\n• 타겟이 명확한 의도적 리퍼럴\n• 신뢰의 서클 (Trust Circle) 기반\n• 구조적 협력으로 기회 배가\n\n최소 4명의 비경쟁 비즈니스가 협력하여 이상적 고객을 공유하고 성장을 가속화',
     'group': '최종'},

    {'id': '4', 'bg': 'dark', 'label': 'WHY POWER TEAM?',
     'title': '무작위 리드 VS 의도적 리퍼럴',
     'body': 'THE OLD WAY                    THE SMART WAY\n불확실한 리드                    타겟 명확 리퍼럴\n낮은 신뢰도                      Trust Circle 기반\n성과 저조                          기회 배가\n\n당신은 어느 쪽을 선택하시겠습니까?',
     'group': '최종'},

    {'id': '5', 'bg': 'light', 'label': 'STATISTICAL EVIDENCE',
     'title': '의도적 리퍼럴은 데이터가 증명합니다',
     'body': 'BNI Korea 85개 챕터 분석 · 4개 지표 모두 p < 0.05\n\n참가율          +0.602\n챕터 규모      +0.355\n멤버 유지율   +0.318\n방문자 활동   +0.231',
     'group': '최종'},

    {'id': '6', 'bg': 'light', 'label': 'DATA COMPARISON',
     'title': '상위 25% vs 하위 25%',
     'body': '의도적 리퍼럴을 실천하는 챕터 vs 그렇지 않은 챕터\n\n+24% 챕터 규모 (44.3명 vs 35.6명)\n+19% 멤버 유지율 (74.4% vs 62.4%)\n+30% 방문자 활동 (1.63명 vs 1.25명)',
     'group': '최종'},

    {'id': '7', 'bg': 'dark', 'label': 'SOLUTION',
     'title': '체계적 리퍼럴 파트너 영입,\n파워팀 툴킷과 함께',
     'body': 'V·C·P 관계 측정 → RPI 점수 산출 → 챕터 건강도 확인 → 지속 성장\n\n회원 평균 ROI: 회원비의 2배 이상',
     'group': '최종'},

    {'id': '8', 'bg': 'light', 'label': 'HOW IT WORKS',
     'title': '레드라이트에서 그린라이트로',
     'body': '리퍼럴 파트너 영입 현황을 한눈에\n\n나의 입력 현황과 챕터 멤버 전원의 현황을 한 화면에서 비교할 수 있습니다.\n누가, 어떤 전문분야의 파트너를 찾고 있는지 서로 확인하고 함께 협력하세요.\n\n🔴 레드라이트 — 전문분야만 입력 (아직 영입 전)\n🟢 그린라이트 — 파트너 성함 입력 (실제 영입 완료)\n\n[대시보드 이미지: 스크린샷 2026-02-27 155022.png]',
     'group': '최종'},

    {'id': '9', 'bg': 'light', 'label': 'HOW TO USE',
     'title': '이렇게 사용하시면 됩니다',
     'body': '복잡하지 않습니다. 3분이면 충분합니다\n\n① 로그인 — powerteam-bnikorea.com 접속 후 간편 로그인\n② 정보 입력 — 파트너 정보 입력, \'수익 단계\' 체크 필수\n③ 제출하기 — 버튼 한 번으로 끝! 챕터 RPI에 즉시 반영',
     'group': '최종'},

    {'id': '10', 'bg': 'dark', 'label': '',
     'title': '무작위에서\n의도적으로,\n지금 시작하세요',
     'body': 'powerteam-bnikorea.com\n\nGivers Gain · BNI Korea',
     'group': '최종'},

    {'id': '11', 'bg': 'light', 'label': 'Q & A',
     'title': '감사합니다',
     'body': '질문이 있으시면 말씀해주세요\npowerteam-bnikorea.com',
     'group': '최종'},

    # === 추가 슬라이드 B1-B8 ===
    {'id': 'B1', 'bg': 'light', 'label': 'GLOBAL CONTEXT',
     'title': '글로벌 파워팀의 철학을 디지털 툴킷으로',
     'body': '1단계: 비즈니스 협업 역량진단\n2단계: 핵심 고객 프로필 매핑\n3단계(오늘의 주제): RPS 리퍼럴 파트너 스코어보드\n\n"Referrals are Intentional, not Random." — BNI Global Power Team',
     'group': '추가'},

    {'id': 'B2', 'bg': 'dark', 'label': 'BNI KOREA POWERTEAM TOOLKIT',
     'title': 'RPS로\n리퍼럴 파트너를 확보하세요',
     'body': 'Referral Partner Scoreboard · 사용 가이드 & 데이터 인사이트\n\n35+ 참여 국가    |    85 활성 챕터    |    0.62 평균 RPI',
     'group': '추가'},

    {'id': 'B3', 'bg': 'light', 'label': 'CORE CONCEPT',
     'title': 'V → C → P',
     'body': '관계의 3단계: 가시성 → 신뢰성 → 수익\n\nVisibility(가시성): "무엇을 하시는지 알아요"\nCredibility(신뢰성): "믿고 추천할 수 있어요"\nProfitability(수익): "비즈니스가 오가고 있어요!"\n\nP 체크 횟수 = RPI 점수 → 챕터 건강도',
     'group': '추가'},

    {'id': 'B4', 'bg': 'light', 'label': 'HOW TO USE',
     'title': '3분이면 끝',
     'body': 'rps-bnikorea.com · BNI Connect 이메일 + 4자리 비밀번호\n\n① 접속 & 로그인\n② 파트너별 VCR 단계 선택\n③ 제출 — 끝!',
     'group': '추가'},

    {'id': 'B5', 'bg': 'light', 'label': 'DATA COMPARISON',
     'title': '상위 25% vs 하위 25%',
     'body': 'RPI ≥ 0.91 vs RPI ≤ 0.22 · p < 0.05\n\n+24% 챕터 규모\n+19% 멤버 유지율\n+30% 방문자 활동',
     'group': '추가'},

    {'id': 'B6', 'bg': 'light', 'label': 'SUCCESS MODELS',
     'title': 'TOP 챕터의 성장 공식',
     'body': 'RPI 2.0 이상 + 참가율 89% 이상\n\nCrown 2.79 (참가율 97.9%)\nSmart 2.63 (참가율 89.3%)\nSynergy 2.35 (참가율 98.7%)\n\n높은 RPI → 참가율↑ → 규모↑ → 유지율↑ → 지속 성장',
     'group': '추가'},

    {'id': 'B7', 'bg': 'dark', 'label': 'KEY INSIGHT',
     'title': 'RPS가 강력한 3가지 이유',
     'body': '01. 관계를 측정 가능하게\n02. 전략적 파트너십 설계\n03. 챕터 전체 선순환 성장',
     'group': '추가'},

    {'id': 'B8', 'bg': 'dark', 'label': '',
     'title': '매주 3분이\n챕터의 미래를 바꿉니다',
     'body': '지금 바로 시작하기\nrps-bnikorea.com\n\nGivers Gain · BNI Korea × TOWN X',
     'group': '추가'},

    # === 초안 APPENDIX A1-A7 ===
    {'id': 'A1', 'bg': 'dark', 'label': 'APPENDIX · ORIGINAL DRAFT',
     'title': '파워팀 툴킷 완벽 가이드',
     'body': 'Givers Gain · BNI KOREA\nRPS 관리 · RPI 연동 · 입력수치 최적화\n10분 핵심 프레젠테이션',
     'group': '초안'},

    {'id': 'A2', 'bg': 'light', 'label': 'WHY POWER TEAM?',
     'title': '무작위 리드 VS 의도적 리퍼럴',
     'body': 'THE OLD WAY: "Leads are Random"\n• 불확실한 무작위 리드\n• 낮은 신뢰와 일방적 관계\n• 낮은 전환율과 성과 저조\n\nTHE SMART WAY: "Referrals are Intentional"\n• 타겟이 명확한 의도적 리퍼럴\n• 신뢰의 서클 (Trust Circle) 기반\n• 구조적 협력으로 기회 배가\n\n파워팀의 힘: 최소 4명의 비경쟁 비즈니스 협력 · ROI: 회원비의 2배 이상',
     'group': '초안'},

    {'id': 'A3', 'bg': 'light', 'label': 'STATISTICAL EVIDENCE',
     'title': 'RPI와 챕터 건강도의 강력한 상관관계',
     'body': '"RPI는 챕터 건강도의 핵심 지표입니다" · 4가지 핵심 지표 (p < 0.05)\n\n참가율 +0.602 (Strong, p<0.001)\n챕터 규모 +0.355 (Moderate, p=0.0012)\n멤버 유지율 +0.318 (Moderate, p=0.0041)\n방문자 활동 +0.231 (Weak, p=0.0392)',
     'group': '초안'},

    {'id': 'A4', 'bg': 'light', 'label': 'DATA COMPARISON',
     'title': 'TOP 25% vs BOTTOM 25%: 데이터로 본 차이',
     'body': 'LOW RPI (하위25%, ≤0.22): 35.6명 / 62.4% / 1.25\nHIGH RPI (상위25%, ≥0.91): 44.3명 / 74.4% / 1.63\n\n+24% 더 큰 규모 / +19% 더 높은 안정성 / +30% 더 많은 기회\n\nHIGH RPI → 규모 확대 → 조직 안정화 → 외부 유입↑ → 지속 성장',
     'group': '초안'},

    {'id': 'A5', 'bg': 'light', 'label': 'STATISTICAL EVIDENCE',
     'title': '통계적 증거: 데이터가 말합니다',
     'body': 'p < 0.05 = 통계적으로 유의미함\n\n+0.602 참가율 (강한 긍정적, p<0.001)\n+0.355 챕터 규모 (중간 긍정적, p=0.0012)\n+0.318 멤버 유지율 (중간 긍정적, p=0.0041)\n+0.231 방문자 활동 (약한 긍정적, p=0.0392)\n\n"RPI는 챕터의 건강도를 측정하는 가장 신뢰할 수 있는 지표입니다."',
     'group': '초안'},

    {'id': 'A6', 'bg': 'light', 'label': 'QUICK START',
     'title': '복잡하지 않습니다. 3분이면 충분합니다',
     'body': '누구나 쉽게 시작하는 3단계 가이드\n\n01. 로그인(1분) — powerteam-bnikorea.com 접속\n02. 정보 입력(1분) — 파트너 정보 입력, 수익 단계 체크 필수\n03. 제출하기(30초) — 버튼 한 번으로 끝! 챕터 RPI에 즉시 반영\n\n당신의 챕터도 TOP 5에 진입할 수 있습니다',
     'group': '초안'},

    {'id': 'A7', 'bg': 'dark', 'label': 'CHAPTER HEALTH',
     'title': '건강한 챕터가 성공합니다',
     'body': 'RPI 높은 챕터 = 건강한 챕터 = 지속 가능한 성장\n\nCrown RPI 2.79 (참가율 97.9%, 45명)\nSmart RPI 2.63 (참가율 89.3%)\nHonors RPI 2.67 (참가율 96.7%)\n\nACTION 1: 지금 파워팀 툴킷 로그인\nACTION 2: 매주 RPS 입력 및 관리\nACTION 3: AI 분석으로 건강도 UP\n\n당신의 챕터도 TOP 25%에 진입할 수 있습니다',
     'group': '초안'},
]

total = len(slides_data)

for idx, sd in enumerate(slides_data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    is_dark = sd['bg'] == 'dark'
    if is_dark:
        set_bg_dark(slide)

    title_color = WHITE if is_dark else TEXT_PRIMARY
    body_color = RGBColor(0x88, 0x88, 0x99) if is_dark else TEXT_SECONDARY
    label_color = RGBColor(0x66, 0x66, 0x77) if is_dark else ACCENT

    # My PowerTeam (top-left)
    add_site_name(slide, dark=is_dark)

    # Group tag (below site name)
    tag = f"[{sd['group']}] #{sd['id']}"
    add_text(slide, Emu(200000), Emu(450000), Emu(3000000), Emu(300000),
             tag, size=9, bold=True, color=ACCENT, align=PP_ALIGN.LEFT)

    # Logo (top-right)
    add_logo(slide, dark=is_dark)

    # Section label
    if sd['label']:
        add_text(slide, Emu(1000000), Emu(1400000), Emu(10192000), Emu(400000),
                 sd['label'], size=11, bold=True, color=label_color)

    # Title
    title_top = Emu(1900000) if sd['label'] else Emu(1600000)
    add_text(slide, Emu(1000000), title_top, Emu(10192000), Emu(900000),
             sd['title'], size=32, bold=True, color=title_color, font_name='Paperlogy')

    # Accent line
    add_accent_line(slide, Emu(3100000))

    # Body
    add_text(slide, Emu(1200000), Emu(3400000), Emu(9792000), Emu(3100000),
             sd['body'], size=14, color=body_color)

    # Copyright (bottom-left)
    add_copyright(slide, dark=is_dark)

    # Page number (bottom-right)
    add_page_num(slide, idx + 1, total, dark=is_dark)

    # Bottom bar
    add_bottom_bar(slide)

output = 'RPS_presentation_review.pptx'
prs.save(output)
print(f'저장 완료: {output}')
print(f'총 슬라이드: {total}장')
