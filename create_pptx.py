"""
RPS (Referral Partner Scoreboard) 프레젠테이션 - PPTX 생성 스크립트
BNI Korea PowerTeam Toolkit
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Colors ──
RED = RGBColor(0xCC, 0x00, 0x00)
DARK = RGBColor(0x1A, 0x1A, 0x2E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0x99, 0x99, 0x99)
MED_GRAY = RGBColor(0x66, 0x66, 0x66)
GREEN = RGBColor(0x00, 0xB8, 0x94)
GOLD = RGBColor(0xF3, 0x9C, 0x12)
LIGHT_RED_BG = RGBColor(0xFF, 0xF0, 0xF0)
NEAR_WHITE = RGBColor(0xF8, 0xF9, 0xFA)
BG_DARK = RGBColor(0x1A, 0x1A, 0x2E)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def add_textbox(slide, left, top, width, height, text, font_size=18,
                color=DARK, bold=False, alignment=PP_ALIGN.LEFT, font_name='맑은 고딕'):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_paragraph(text_frame, text, font_size=14, color=DARK, bold=False,
                  alignment=PP_ALIGN.LEFT, font_name='맑은 고딕', space_before=0):
    p = text_frame.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    if space_before:
        p.space_before = Pt(space_before)
    return p

def set_slide_bg(slide, color):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_red_bar(slide):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(7.4), Inches(13.333), Inches(0.1)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RED
    shape.line.fill.background()

def add_rounded_rect(slide, left, top, width, height, fill_color, text='',
                     font_size=14, font_color=WHITE, bold=False, alignment=PP_ALIGN.CENTER):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = font_color
        p.font.bold = bold
        p.font.name = '맑은 고딕'
        p.alignment = alignment
        tf.paragraphs[0].space_before = Pt(4)
    return shape

def add_circle(slide, left, top, size, fill_color, text='', font_size=20, font_color=WHITE):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(left), Inches(top), Inches(size), Inches(size)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = font_color
        p.font.bold = True
        p.font.name = '맑은 고딕'
        p.alignment = PP_ALIGN.CENTER
    return shape

# ══════════════════════════════════════════════════════════════
# SLIDE 1: TITLE
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
set_slide_bg(slide, BG_DARK)

add_textbox(slide, 0, 0.5, 13.333, 0.5,
            'BNI KOREA × TOWN X  |  POWERTEAM TOOLKIT',
            font_size=14, color=RGBColor(0xFF, 0x6B, 0x6B), bold=True, alignment=PP_ALIGN.CENTER)

add_textbox(slide, 1, 1.5, 11.333, 1.5,
            'RPS로\n리퍼럴 파트너를\n확보하세요',
            font_size=52, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# Red line
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.9), Inches(3.7), Inches(1.5), Inches(0.06))
shape.fill.solid()
shape.fill.fore_color.rgb = RED
shape.line.fill.background()

add_textbox(slide, 2, 3.9, 9.333, 0.5,
            'Referral Partner Scoreboard  |  사용 가이드 & 데이터 인사이트',
            font_size=16, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# Bottom stats
stats = [('35+', 'Global Power Team\n참여 국가'), ('85', 'BNI Korea\n활성 챕터'), ('0.62', '현재 BNI Korea\n평균 RPI')]
for i, (num, label) in enumerate(stats):
    x = 3.8 + i * 2.2
    c = RED if i == 2 else WHITE
    add_textbox(slide, x, 4.8, 1.8, 0.6, num, font_size=30, color=c, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x, 5.4, 1.8, 0.6, label, font_size=10, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

add_red_bar(slide)

# ══════════════════════════════════════════════════════════════
# SLIDE 2: GLOBAL CONTEXT
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])

add_textbox(slide, 0.7, 0.4, 6, 0.3, 'GLOBAL CONTEXT', font_size=14, color=RED, bold=True)
add_textbox(slide, 0.7, 0.7, 10, 0.7, '글로벌 검증 시스템, 한국형 디지털 혁신', font_size=36, color=DARK, bold=True)
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.7), Inches(1.4), Inches(1.2), Inches(0.06))
shape.fill.solid(); shape.fill.fore_color.rgb = RED; shape.line.fill.background()

# Left - Global Power Team info
box = add_rounded_rect(slide, 0.7, 1.8, 5.8, 2.2, BG_DARK)
tf = box.text_frame; tf.word_wrap = True
tf.paragraphs[0].text = 'BNI GLOBAL POWER TEAM (Since 2010)'
tf.paragraphs[0].font.size = Pt(12); tf.paragraphs[0].font.color.rgb = LIGHT_GRAY
tf.paragraphs[0].font.name = '맑은 고딕'
for txt, clr in [
    ('35+ 참여 국가  |  4가지 핵심 가치', WHITE),
    ('', WHITE),
    ('🌍 글로벌 네트워킹 — 35개국 이상 관계 구축', RGBColor(0xCC,0xCC,0xCC)),
    ('📈 비즈니스 성장 — 고품질 소개 & 신규 시장 진출', RGBColor(0xCC,0xCC,0xCC)),
    ('🎓 전문성 개발 — 국제 비즈니스 역량 강화', RGBColor(0xCC,0xCC,0xCC)),
    ('🤝 국경 간 협력 — 파트너십 & 지식 공유', RGBColor(0xCC,0xCC,0xCC)),
]:
    add_paragraph(tf, txt, font_size=12 if clr == RGBColor(0xCC,0xCC,0xCC) else 16,
                  color=clr, bold=(clr == WHITE))

# Quote
add_textbox(slide, 0.7, 4.2, 5.8, 0.6,
            '"Referrals are Intentional, not Random."\n— BNI Global Power Team',
            font_size=13, color=MED_GRAY)

# Right - Korea Innovation
box2 = add_rounded_rect(slide, 7, 1.8, 5.6, 3.6, WHITE)
box2.line.color.rgb = RED; box2.line.width = Pt(2)
tf2 = box2.text_frame; tf2.word_wrap = True
tf2.paragraphs[0].text = ''; tf2.paragraphs[0].font.size = Pt(4)

add_paragraph(tf2, '🇰🇷 한국의 혁신', font_size=16, color=RED, bold=True, alignment=PP_ALIGN.CENTER)
add_paragraph(tf2, '', font_size=6)
add_paragraph(tf2, '글로벌 파워팀의 철학을 디지털 툴킷으로 발전:', font_size=13, color=MED_GRAY)
add_paragraph(tf2, '', font_size=6)
add_paragraph(tf2, '1단계  비즈니스 협업 역량진단 (10분)', font_size=14, color=DARK, bold=True)
add_paragraph(tf2, '2단계  핵심 고객 프로필 매핑 (5분)', font_size=14, color=DARK, bold=True)
add_paragraph(tf2, '3단계  RPS 리퍼럴 파트너 스코어보드 (실시간)', font_size=14, color=RED, bold=True)
add_paragraph(tf2, '', font_size=8)
add_paragraph(tf2, '▶ 오늘의 주제 → 3단계 RPS', font_size=14, color=RED, bold=True, alignment=PP_ALIGN.CENTER)

# Bottom banner
add_rounded_rect(slide, 0.7, 5.7, 11.9, 0.6, BG_DARK,
    '💡 글로벌 파워팀의 "의도적 리퍼럴" 철학을 데이터 기반으로 실현하는 도구 = RPS',
    font_size=13, font_color=WHITE)

add_red_bar(slide)

# ══════════════════════════════════════════════════════════════
# SLIDE 3: PAIN POINT
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])

add_textbox(slide, 0, 0.4, 13.333, 0.3, 'THE PROBLEM', font_size=14, color=RED, bold=True, alignment=PP_ALIGN.CENTER)
add_textbox(slide, 0, 0.8, 13.333, 0.9, '리퍼럴 파트너,\n어떻게 관리하고 계신가요?', font_size=36, color=DARK, bold=True, alignment=PP_ALIGN.CENTER)

# Left box - Old Way
box_old = add_rounded_rect(slide, 0.8, 2.2, 5.2, 4, RGBColor(0xF5, 0xF5, 0xF5))
tf = box_old.text_frame; tf.word_wrap = True
tf.paragraphs[0].text = '😰 지금의 현실'; tf.paragraphs[0].font.size = Pt(22); tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.color.rgb = DARK; tf.paragraphs[0].font.name = '맑은 고딕'; tf.paragraphs[0].alignment = PP_ALIGN.CENTER
add_paragraph(tf, '"Leads are Random"', font_size=11, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
add_paragraph(tf, '', font_size=8)
for txt in ['❌ 누구와 어떤 단계인지 기억이...', '❌ 엑셀? 메모장? 머릿속?',
            '❌ 챕터 RPI가 왜 낮은지 모름', '❌ 파트너십 전략? 감으로 한다']:
    add_paragraph(tf, txt, font_size=15, color=LIGHT_GRAY, space_before=8)

# Arrow
add_textbox(slide, 6.1, 3.5, 1.1, 0.6, '→', font_size=40, color=DARK, bold=True, alignment=PP_ALIGN.CENTER)

# Right box - RPS
box_new = add_rounded_rect(slide, 7.3, 2.2, 5.2, 4, WHITE)
box_new.line.color.rgb = RED; box_new.line.width = Pt(3)
tf = box_new.text_frame; tf.word_wrap = True
tf.paragraphs[0].text = '🎯 RPS가 해결합니다'; tf.paragraphs[0].font.size = Pt(22); tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.color.rgb = RED; tf.paragraphs[0].font.name = '맑은 고딕'; tf.paragraphs[0].alignment = PP_ALIGN.CENTER
add_paragraph(tf, '"Referrals are Intentional"', font_size=11, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
add_paragraph(tf, '', font_size=8)
for txt in ['✅ 파트너별 VCR 단계 한눈에', '✅ 웹 기반 실시간 자동 관리',
            '✅ 챕터 RPI에 즉시 반영', '✅ 전국 85개 챕터 비교 분석']:
    add_paragraph(tf, txt, font_size=15, color=DARK, bold=True, space_before=8)

add_red_bar(slide)

# ══════════════════════════════════════════════════════════════
# SLIDE 4: VCR CONCEPT
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])

add_textbox(slide, 0.7, 0.4, 6, 0.3, 'CORE CONCEPT', font_size=14, color=RED, bold=True)
add_textbox(slide, 0.7, 0.7, 10, 0.6, 'RPS의 핵심: V → C → P', font_size=36, color=DARK, bold=True)
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.7), Inches(1.3), Inches(1.2), Inches(0.06))
shape.fill.solid(); shape.fill.fore_color.rgb = RED; shape.line.fill.background()
add_textbox(slide, 0.7, 1.5, 10, 0.4, '글로벌 파워팀이 검증한 리퍼럴 파트너십의 3단계', font_size=14, color=LIGHT_GRAY)

# V Card
cards_data = [
    ('👁️ Visibility', '가시성', GREEN, '"이 분이 무엇을 하시는지 알아요"', RGBColor(0xE6, 0xFF, 0xF5)),
    ('🤝 Credibility', '신뢰성', GOLD, '"이 분을 믿고 추천할 수 있어요"', RGBColor(0xFF, 0xF8, 0xE1)),
    ('💰 Profitability', '수익 발생', RED, '"실제 비즈니스가 오가고 있어요!"', RGBColor(0xFF, 0xF0, 0xF0)),
]
for i, (title, sub, color, quote_text, bg) in enumerate(cards_data):
    x = 0.8 + i * 4.2
    if i == 2:
        box = add_rounded_rect(slide, x, 2.1, 3.8, 3.0, RED)
        tf = box.text_frame; tf.word_wrap = True
        tf.paragraphs[0].text = title; tf.paragraphs[0].font.size = Pt(22)
        tf.paragraphs[0].font.bold = True; tf.paragraphs[0].font.color.rgb = WHITE; tf.paragraphs[0].font.name = '맑은 고딕'
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        add_paragraph(tf, sub, font_size=14, color=RGBColor(0xFF,0xCC,0xCC), alignment=PP_ALIGN.CENTER)
        add_paragraph(tf, '', font_size=8)
        add_paragraph(tf, quote_text, font_size=13, color=WHITE, alignment=PP_ALIGN.CENTER)
    else:
        box = add_rounded_rect(slide, x, 2.1, 3.8, 3.0, WHITE)
        # Top border line
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(2.1), Inches(3.8), Inches(0.08))
        line.fill.solid(); line.fill.fore_color.rgb = color; line.line.fill.background()
        tf = box.text_frame; tf.word_wrap = True
        tf.paragraphs[0].text = title; tf.paragraphs[0].font.size = Pt(22)
        tf.paragraphs[0].font.bold = True; tf.paragraphs[0].font.color.rgb = DARK; tf.paragraphs[0].font.name = '맑은 고딕'
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        add_paragraph(tf, sub, font_size=14, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
        add_paragraph(tf, '', font_size=8)
        add_paragraph(tf, quote_text, font_size=13, color=MED_GRAY, alignment=PP_ALIGN.CENTER)

# Highlight box
add_rounded_rect(slide, 0.8, 5.4, 11.7, 0.7, LIGHT_RED_BG,
    '핵심 공식: RPS에서 \'수익 단계(P)\' 체크 횟수가 곧 RPI 점수에 반영됩니다. P 체크 ↑ → RPI ↑ → 챕터 건강도 UP',
    font_size=13, font_color=DARK, alignment=PP_ALIGN.LEFT)

add_red_bar(slide)

# ══════════════════════════════════════════════════════════════
# SLIDE 5: HOW TO USE
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])

add_textbox(slide, 0.7, 0.4, 6, 0.3, 'HOW TO USE', font_size=14, color=RED, bold=True)
add_textbox(slide, 0.7, 0.7, 10, 0.6, '3분이면 끝. 정말 간단합니다', font_size=36, color=DARK, bold=True)
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.7), Inches(1.3), Inches(1.2), Inches(0.06))
shape.fill.solid(); shape.fill.fore_color.rgb = RED; shape.line.fill.background()

# Try to add phone image
phone_img = os.path.join(os.path.dirname(__file__), '스크린샷_2026-02-27_154050-removebg-preview.png')
if os.path.exists(phone_img):
    slide.shapes.add_picture(phone_img, Inches(0.8), Inches(1.7), height=Inches(4.8))

# Steps
steps = [
    ('1', '접속 & 로그인', 'rps-bnikorea.com 접속\nBNI Connect 이메일 + 4자리 비밀번호'),
    ('2', '파트너별 VCR 단계 선택', '각 파트너와의 관계 단계를 선택하고\n\'수익 단계\' 반드시 체크!'),
    ('3', '제출 — 끝!', '버튼 한 번이면 완료\n챕터 RPI에 즉시 반영'),
]
for i, (num, title, desc) in enumerate(steps):
    y = 1.8 + i * 1.5
    add_circle(slide, 6.5, y, 0.65, RED, num, font_size=22)
    add_textbox(slide, 7.4, y, 5, 0.4, title, font_size=20, color=DARK, bold=True)
    add_textbox(slide, 7.4, y + 0.45, 5, 0.7, desc, font_size=13, color=LIGHT_GRAY)

# Time box
add_rounded_rect(slide, 6.5, 5.8, 6, 0.65, BG_DARK,
    '⏱️ 매주 3분 투자 → 챕터 성장의 시작', font_size=16, font_color=WHITE, bold=True)

add_red_bar(slide)

# ══════════════════════════════════════════════════════════════
# SLIDE 6: DASHBOARD
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])

add_textbox(slide, 0.7, 0.4, 6, 0.3, 'RPS DASHBOARD', font_size=14, color=RED, bold=True)
add_textbox(slide, 0.7, 0.7, 10, 0.6, '숫자로 보는 우리 챕터의 위치', font_size=36, color=DARK, bold=True)
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.7), Inches(1.3), Inches(1.2), Inches(0.06))
shape.fill.solid(); shape.fill.fore_color.rgb = RED; shape.line.fill.background()

# Left table - Regional
add_textbox(slide, 0.8, 1.6, 5.5, 0.4, '📊 전국 17개 지역 RPI 현황', font_size=16, color=DARK, bold=True)

region_data = [
    ('지역', 'RPI', '참가율'),
    ('Busan 부산', '0.74', '98.9%'),
    ('Seocho 서초', '0.91', '85.6%'),
    ('Goyang 고양', '0.32', '91.5%'),
    ('Songpa 송파', '0.86', '53.9%'),
    ('BNI Korea 전체', '0.62', '71.0%'),
]
from pptx.util import Inches as In
table_shape = slide.shapes.add_table(len(region_data), 3, Inches(0.8), Inches(2.1), Inches(5.5), Inches(3.0))
table = table_shape.table
for r, row_data in enumerate(region_data):
    for c, val in enumerate(row_data):
        cell = table.cell(r, c)
        cell.text = val
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(13)
            p.font.name = '맑은 고딕'
            if r == 0:
                p.font.color.rgb = WHITE
                p.font.bold = True
                cell.fill.solid()
                cell.fill.fore_color.rgb = RED
            elif r == len(region_data) - 1:
                p.font.color.rgb = RED
                p.font.bold = True
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_RED_BG
            else:
                p.font.color.rgb = DARK

# Right table - Chapter ranking
add_textbox(slide, 7, 1.6, 5.5, 0.4, '🏆 챕터 RPI TOP 5 vs BOTTOM', font_size=16, color=DARK, bold=True)

chapter_data = [
    ('챕터', 'RPI', '참가율'),
    ('🥇 Smart', '2.63', '89.3%'),
    ('🥈 Synergy', '2.35', '98.7%'),
    ('🥉 Winners', '1.42', '94.4%'),
    ('GROW', '0.00', '30.0%'),
    ('Stoneworks', '0.00', '19.0%'),
]
table_shape2 = slide.shapes.add_table(len(chapter_data), 3, Inches(7), Inches(2.1), Inches(5.5), Inches(3.0))
table2 = table_shape2.table
for r, row_data in enumerate(chapter_data):
    for c, val in enumerate(row_data):
        cell = table2.cell(r, c)
        cell.text = val
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(13)
            p.font.name = '맑은 고딕'
            if r == 0:
                p.font.color.rgb = WHITE; p.font.bold = True
                cell.fill.solid(); cell.fill.fore_color.rgb = RED
            elif r <= 2:
                p.font.color.rgb = RED if c == 1 else DARK
                p.font.bold = True
                cell.fill.solid(); cell.fill.fore_color.rgb = LIGHT_RED_BG
            elif r >= 4:
                p.font.color.rgb = LIGHT_GRAY
            else:
                p.font.color.rgb = DARK

add_textbox(slide, 7, 5.2, 5.5, 0.4, '💡 RPI 0.00 챕터 vs 2.63 챕터 = 참가율 70%p 차이',
            font_size=12, color=RED, bold=True, alignment=PP_ALIGN.CENTER)

add_rounded_rect(slide, 0.8, 5.7, 11.7, 0.6, BG_DARK,
    '📊 전국 17개 지역 · 85개 챕터의 RPI & 참가율을 실시간 비교 → 우리 챕터의 위치를 객관적으로 파악',
    font_size=12, font_color=WHITE)

add_red_bar(slide)

# ══════════════════════════════════════════════════════════════
# SLIDE 7: STATISTICAL EVIDENCE
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, NEAR_WHITE)

add_textbox(slide, 0.7, 0.4, 6, 0.3, 'STATISTICAL EVIDENCE', font_size=14, color=RED, bold=True)
add_textbox(slide, 0.7, 0.7, 10, 0.6, '데이터가 증명합니다: RPI의 힘', font_size=36, color=DARK, bold=True)
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.7), Inches(1.3), Inches(1.2), Inches(0.06))
shape.fill.solid(); shape.fill.fore_color.rgb = RED; shape.line.fill.background()
add_textbox(slide, 0.7, 1.5, 12, 0.3, 'BNI Korea 85개 챕터 분석  |  4개 지표 모두 p < 0.05 통계적 유의미',
            font_size=13, color=LIGHT_GRAY)

# Correlation bars
corr_data = [
    ('👥 참가율 (Participation)', '+0.602', 'p < 0.001', '강한 긍정', 60.2),
    ('📐 챕터 규모 (Size)', '+0.355', 'p = 0.001', '중간 긍정', 35.5),
    ('🔄 멤버 유지율 (Retention)', '+0.318', 'p = 0.004', '중간 긍정', 31.8),
    ('🚪 방문자 활동 (Visitors)', '+0.231', 'p = 0.039', '약한 긍정', 23.1),
]

for i, (label, coeff, p_val, strength, width_pct) in enumerate(corr_data):
    y = 2.1 + i * 1.0
    add_textbox(slide, 0.8, y, 5, 0.35, label, font_size=15, color=DARK, bold=True)
    add_textbox(slide, 7, y, 2, 0.35, coeff, font_size=22, color=GREEN, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, 9.2, y, 1.5, 0.35, p_val, font_size=11, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, 10.8, y, 1.5, 0.35, strength, font_size=11, color=RED, bold=True, alignment=PP_ALIGN.CENTER)

    # Bar background
    bar_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.8), Inches(y + 0.4), Inches(11.5), Inches(0.22))
    bar_bg.fill.solid(); bar_bg.fill.fore_color.rgb = RGBColor(0xE0, 0xE0, 0xE0); bar_bg.line.fill.background()

    # Bar fill
    bar_fill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.8), Inches(y + 0.4), Inches(11.5 * width_pct / 100), Inches(0.22))
    bar_fill.fill.solid(); bar_fill.fill.fore_color.rgb = GREEN; bar_fill.line.fill.background()

# Scale
add_textbox(slide, 0.8, 6.2, 11.5, 0.3,
    'No Correlation (0)          Weak (+0.2)          Moderate (+0.4)          Strong (+0.6)          Perfect (+1.0)',
    font_size=10, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

add_rounded_rect(slide, 0.8, 6.5, 11.7, 0.55, GREEN,
    '💡 결론: RPI는 챕터 건강도의 핵심 지표입니다. 4개 지표 모두 통계적으로 유의미한 긍정적 상관관계 확인',
    font_size=12, font_color=WHITE, bold=True)

add_red_bar(slide)

# ══════════════════════════════════════════════════════════════
# SLIDE 8: TOP 25% vs BOTTOM 25%
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])

add_textbox(slide, 0, 0.3, 13.333, 0.3, 'DATA COMPARISON', font_size=14, color=RED, bold=True, alignment=PP_ALIGN.CENTER)
add_textbox(slide, 0, 0.6, 13.333, 0.6, '상위 25% vs 하위 25%: 숫자의 격차', font_size=36, color=DARK, bold=True, alignment=PP_ALIGN.CENTER)
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.1), Inches(1.2), Inches(1.2), Inches(0.06))
shape.fill.solid(); shape.fill.fore_color.rgb = RED; shape.line.fill.background()

# Left - Low RPI
box_low = add_rounded_rect(slide, 0.8, 1.6, 5.6, 2.6, RGBColor(0xF0, 0xF0, 0xF0))
tf = box_low.text_frame; tf.word_wrap = True
tf.paragraphs[0].text = '📉 LOW RPI 챕터 (하위 25%)'; tf.paragraphs[0].font.size = Pt(16)
tf.paragraphs[0].font.bold = True; tf.paragraphs[0].font.color.rgb = DARK; tf.paragraphs[0].font.name = '맑은 고딕'
add_paragraph(tf, '평균 RPI ≤ 0.22', font_size=11, color=LIGHT_GRAY)
add_paragraph(tf, '', font_size=6)
add_paragraph(tf, '평균 챕터 규모          35.6명', font_size=15, color=DARK, bold=True)
add_paragraph(tf, '멤버 유지율               62.4%', font_size=15, color=DARK, bold=True, space_before=6)
add_paragraph(tf, '방문자/멤버               1.25명', font_size=15, color=DARK, bold=True, space_before=6)

# Right - High RPI
box_high = add_rounded_rect(slide, 6.9, 1.6, 5.6, 2.6, RED)
tf = box_high.text_frame; tf.word_wrap = True
tf.paragraphs[0].text = '📈 HIGH RPI 챕터 (상위 25%)'; tf.paragraphs[0].font.size = Pt(16)
tf.paragraphs[0].font.bold = True; tf.paragraphs[0].font.color.rgb = WHITE; tf.paragraphs[0].font.name = '맑은 고딕'
add_paragraph(tf, '평균 RPI ≥ 0.91', font_size=11, color=RGBColor(0xFF, 0xCC, 0xCC))
add_paragraph(tf, '', font_size=6)
add_paragraph(tf, '평균 챕터 규모          44.3명', font_size=15, color=WHITE, bold=True)
add_paragraph(tf, '멤버 유지율               74.4%', font_size=15, color=WHITE, bold=True, space_before=6)
add_paragraph(tf, '방문자/멤버               1.63명', font_size=15, color=WHITE, bold=True, space_before=6)

# Impact numbers
impact_data = [
    ('+24%', '챕터 규모', '+8.7명 더 큰 규모', RED),
    ('+19%', '멤버 유지율', '+12%p 더 높은 안정성', GREEN),
    ('+30%', '방문자 활동', '+0.38명/멤버 더 많은 기회', GOLD),
]
for i, (num, title, desc, color) in enumerate(impact_data):
    x = 1.2 + i * 4.0
    box = add_rounded_rect(slide, x, 4.6, 3.4, 1.5, WHITE)
    tf = box.text_frame; tf.word_wrap = True
    tf.paragraphs[0].text = num; tf.paragraphs[0].font.size = Pt(36)
    tf.paragraphs[0].font.bold = True; tf.paragraphs[0].font.color.rgb = color
    tf.paragraphs[0].font.name = '맑은 고딕'; tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    add_paragraph(tf, title, font_size=14, color=DARK, bold=True, alignment=PP_ALIGN.CENTER)
    add_paragraph(tf, desc, font_size=10, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

add_textbox(slide, 0, 6.3, 13.333, 0.3, '모든 차이는 통계적으로 유의미합니다 (p < 0.05)',
            font_size=11, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

add_red_bar(slide)

# ══════════════════════════════════════════════════════════════
# SLIDE 9: SUCCESS MODELS
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])

add_textbox(slide, 0.7, 0.4, 6, 0.3, 'SUCCESS MODELS', font_size=14, color=RED, bold=True)
add_textbox(slide, 0.7, 0.7, 10, 0.6, 'TOP 챕터가 증명하는 성장 공식', font_size=36, color=DARK, bold=True)
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.7), Inches(1.3), Inches(1.2), Inches(0.06))
shape.fill.solid(); shape.fill.fore_color.rgb = RED; shape.line.fill.background()

success_data = [
    ('Crown 챕터', 'BEST', '2.79', '97.9%', '챕터 규모 45명', RED),
    ('Smart 챕터', 'HIGH RET.', '2.63', '89.3%', '높은 멤버 유지율', GOLD),
    ('Synergy 챕터', 'CONSISTENT', '2.35', '98.7%', '참가율 최상위', GREEN),
]
for i, (name, badge, rpi, part, desc, color) in enumerate(success_data):
    x = 0.8 + i * 4.2
    box = add_rounded_rect(slide, x, 1.7, 3.8, 2.8, WHITE)
    # Top border
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(1.7), Inches(3.8), Inches(0.06))
    line.fill.solid(); line.fill.fore_color.rgb = color; line.line.fill.background()

    tf = box.text_frame; tf.word_wrap = True
    tf.paragraphs[0].text = ''; tf.paragraphs[0].font.size = Pt(4)
    add_paragraph(tf, f'[{badge}]', font_size=10, color=color, bold=True, alignment=PP_ALIGN.RIGHT)
    add_paragraph(tf, name, font_size=20, color=DARK, bold=True, alignment=PP_ALIGN.CENTER)
    add_paragraph(tf, '', font_size=4)
    add_paragraph(tf, f'RPI  {rpi}', font_size=22, color=RED, bold=True, alignment=PP_ALIGN.CENTER)
    add_paragraph(tf, f'참가율  {part}', font_size=18, color=GREEN, bold=True, alignment=PP_ALIGN.CENTER)
    add_paragraph(tf, desc, font_size=11, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER, space_before=6)

# Growth flow
flow_items = ['높은 RPI', '→', '높은 참가율', '→', '큰 챕터', '→', '높은 유지율', '→', '많은 방문자', '→', '지속 성장']
add_textbox(slide, 0.5, 4.7, 12.3, 0.5,
    '  '.join(flow_items),
    font_size=14, color=DARK, bold=True, alignment=PP_ALIGN.CENTER)

add_rounded_rect(slide, 0.8, 5.4, 11.7, 0.6, BG_DARK,
    '📊 공통점: TOP 챕터들은 모두 RPI 2.0 이상 + 참가율 89% 이상 → RPS 적극 활용의 결과',
    font_size=12, font_color=WHITE, bold=True)

add_red_bar(slide)

# ══════════════════════════════════════════════════════════════
# SLIDE 10: KEY INSIGHT
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_textbox(slide, 0, 0.3, 13.333, 0.3, 'KEY INSIGHT', font_size=14, color=RGBColor(0xFF, 0x6B, 0x6B), bold=True, alignment=PP_ALIGN.CENTER)
add_textbox(slide, 0, 0.7, 13.333, 0.8, 'RPS가 리퍼럴 파트너 확보에\n강력한 3가지 이유', font_size=32, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

insights = [
    ('🔍', '관계를 \'측정 가능\'하게', 'VCR 시각화로\nV→C→P 전환율 추적\n현재 참가율 71%\n→ 데이터가 다음 액션을 알려줍니다'),
    ('🎯', '전략적 파트너십 설계', 'V단계 → 1:1 미팅 증가\nC단계 → 리퍼럴 요청\nP단계 → 수익 확대\n→ 감 아닌 데이터로 전략을'),
    ('📈', '챕터 전체 선순환 성장', '전원 입력 → RPI +24% 규모↑\n높은 RPI → 유지율 +19%↑\n안정된 챕터 → 방문자 +30%↑\n→ 선순환의 시작점 = RPS'),
]
for i, (icon, title, desc) in enumerate(insights):
    x = 0.8 + i * 4.2
    box = add_rounded_rect(slide, x, 1.9, 3.8, 3.5, WHITE)
    tf = box.text_frame; tf.word_wrap = True
    tf.paragraphs[0].text = icon; tf.paragraphs[0].font.size = Pt(28)
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    add_paragraph(tf, title, font_size=16, color=DARK, bold=True, alignment=PP_ALIGN.CENTER)
    add_paragraph(tf, '', font_size=4)
    for line in desc.split('\n'):
        add_paragraph(tf, line, font_size=12, color=MED_GRAY, alignment=PP_ALIGN.LEFT, space_before=3)

add_rounded_rect(slide, 0.8, 5.7, 11.7, 0.6, RGBColor(0x2A, 0x2A, 0x3E),
    '🌍 BNI Global Power Team: "Cross-border Collaboration & High-quality Referrals" → RPS는 이 철학의 한국형 디지털 구현체',
    font_size=11, font_color=LIGHT_GRAY)

add_red_bar(slide)

# ══════════════════════════════════════════════════════════════
# SLIDE 11: ACTION PLAN
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])

add_textbox(slide, 0, 0.3, 13.333, 0.3, 'ACTION PLAN', font_size=14, color=RED, bold=True, alignment=PP_ALIGN.CENTER)
add_textbox(slide, 0, 0.7, 13.333, 0.6, '오늘부터 바로 실행하세요', font_size=36, color=DARK, bold=True, alignment=PP_ALIGN.CENTER)
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.1), Inches(1.2), Inches(1.2), Inches(0.06))
shape.fill.solid(); shape.fill.fore_color.rgb = RED; shape.line.fill.background()

actions = [
    ('🔑', '오늘', 'RPS에 로그인하고\n파트너 정보를 입력하세요', 'rps-bnikorea.com', RED),
    ('📅', '매주 3분', 'VCR 단계 업데이트 +\n\'수익 단계\' 체크 잊지 마세요', '⏱️ 3분 = RPI 반영', GOLD),
    ('🏆', '목표: TOP 25%', '챕터 전원 참여 →\nRPI ≥ 0.91 달성!', '+24% 규모·+19% 유지·+30% 기회', GREEN),
]
for i, (icon, title, desc, sub, color) in enumerate(actions):
    x = 0.8 + i * 4.2
    box = add_rounded_rect(slide, x, 1.6, 3.8, 3.4, BG_DARK)
    add_circle(slide, x + 1.4, 1.9, 0.7, color, icon, font_size=20)
    tf = box.text_frame; tf.word_wrap = True
    tf.paragraphs[0].text = ''; tf.paragraphs[0].font.size = Pt(4)
    add_paragraph(tf, '', font_size=28)
    add_paragraph(tf, title, font_size=18, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_paragraph(tf, '', font_size=4)
    for line in desc.split('\n'):
        add_paragraph(tf, line, font_size=13, color=RGBColor(0xCC,0xCC,0xCC), alignment=PP_ALIGN.CENTER)
    add_paragraph(tf, '', font_size=4)
    add_paragraph(tf, sub, font_size=10, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

add_rounded_rect(slide, 0.8, 5.3, 11.7, 0.7, LIGHT_RED_BG,
    '⚠️ 반드시 기억: 수익 단계(Profitability) 체크가 RPI의 핵심입니다. 파트너와 실제 비즈니스가 오가면 반드시 체크하세요!',
    font_size=13, font_color=DARK, bold=True, alignment=PP_ALIGN.LEFT)

add_red_bar(slide)

# ══════════════════════════════════════════════════════════════
# SLIDE 12: CLOSING
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_textbox(slide, 0, 0.5, 13.333, 0.3, 'YOUR CHAPTER, YOUR GROWTH',
            font_size=13, color=RGBColor(0xFF, 0x6B, 0x6B), bold=True, alignment=PP_ALIGN.CENTER)

add_textbox(slide, 0, 1.0, 13.333, 1.2, '매주 3분의 투자가\n챕터의 미래를 바꿉니다',
            font_size=44, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# Stats row
closing_stats = [
    ('3분', '매주 입력', RED),
    ('+24%', '챕터 규모', GREEN),
    ('+19%', '멤버 유지', GOLD),
    ('+30%', '비즈니스 기회', RGBColor(0xFF, 0x6B, 0x6B)),
]
for i, (num, label, color) in enumerate(closing_stats):
    x = 2.2 + i * 2.4
    add_textbox(slide, x, 3.2, 2, 0.7, num, font_size=38, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x, 3.9, 2, 0.4, label, font_size=12, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

add_rounded_rect(slide, 4.5, 4.8, 4.3, 0.7, RED,
    '지금 바로 시작하기', font_size=20, font_color=WHITE, bold=True)

add_textbox(slide, 0, 5.7, 13.333, 0.3,
    'powerteam-bnikorea.com  |  rps-bnikorea.com',
    font_size=14, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

add_textbox(slide, 0, 6.5, 13.333, 0.3,
    'Givers Gain  |  BNI Korea × TOWN X  |  BNI Global Power Team',
    font_size=10, color=RGBColor(0x55, 0x55, 0x55), alignment=PP_ALIGN.CENTER)

add_red_bar(slide)

# ══════════════════════════════════════════════════════════════
# SLIDE 13: Q&A
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])

add_textbox(slide, 0, 0.8, 13.333, 0.3, 'Q & A', font_size=14, color=RED, bold=True, alignment=PP_ALIGN.CENTER)
add_textbox(slide, 0, 1.5, 13.333, 1.0, '감사합니다', font_size=56, color=DARK, bold=True, alignment=PP_ALIGN.CENTER)

shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.1), Inches(2.7), Inches(1.2), Inches(0.06))
shape.fill.solid(); shape.fill.fore_color.rgb = RED; shape.line.fill.background()

add_textbox(slide, 0, 3.0, 13.333, 0.5, '질문이 있으시면 말씀해주세요',
            font_size=16, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

links = [
    ('🌐 메인 사이트', 'powerteam-bnikorea.com'),
    ('📊 RPS 스코어보드', 'rps-bnikorea.com'),
    ('🌍 글로벌 파워팀', 'bniglobalpowerteam.com'),
]
for i, (title, url) in enumerate(links):
    x = 2.2 + i * 3.2
    box = add_rounded_rect(slide, x, 3.8, 2.8, 1.2, WHITE)
    box.line.color.rgb = RGBColor(0xEE, 0xEE, 0xEE); box.line.width = Pt(2)
    tf = box.text_frame; tf.word_wrap = True
    tf.paragraphs[0].text = title; tf.paragraphs[0].font.size = Pt(14)
    tf.paragraphs[0].font.bold = True; tf.paragraphs[0].font.color.rgb = DARK
    tf.paragraphs[0].font.name = '맑은 고딕'; tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    add_paragraph(tf, url, font_size=11, color=RED, bold=True, alignment=PP_ALIGN.CENTER, space_before=6)

add_textbox(slide, 0, 6.2, 13.333, 0.3,
    'Givers Gain  |  BNI Korea PowerTeam Toolkit  |  BNI Global Power Team',
    font_size=10, color=RGBColor(0xDD, 0xDD, 0xDD), alignment=PP_ALIGN.CENTER)

add_red_bar(slide)

# ── Save ──
output_path = os.path.join(os.path.dirname(__file__), 'RPS_presentation.pptx')
prs.save(output_path)
print(f'PPTX created: {output_path}')
print(f'Total slides: {len(prs.slides)}')
